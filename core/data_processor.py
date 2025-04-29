"""
文档处理器模块 - 论文纠错系统

本模块实现了多种文档格式的加载、解析和处理功能。支持以下文件格式：
- Word文档(.doc, .docx, .docm, .dotm)
- PDF文档(.pdf)
- 纯文本文档(.txt, .text)
- CSV文档(.csv)
- Excel文档(.xls, .xlsx, .xlsm)
- Markdown文档(.md, .markdown)
- PowerPoint文档(.ppt, .pptx, .pptm)

主要功能包括：
- 文档加载与文本提取
- 段落分割与合并
- 文本预处理
- 文档处理器工厂类

@author: Biubush
@date: 2025
"""
import os
import subprocess
import tempfile
import shutil
import time
import re
import abc
import platform
import PyPDF2
import docx2txt
import chardet
import csv
import openpyxl
import markdown
import html2text
import pptx
import sys
from docx import Document

# 仅在Windows平台导入Windows特有库
IS_WINDOWS = platform.system() == "Windows"
if IS_WINDOWS:
    import pythoncom
    from win32com import client as wc

# 导入配置文件
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from config import MAX_CHARS_PER_PARAGRAPH

class DocumentProcessor(abc.ABC):
    """
    文档处理器的抽象基类，定义了所有文档处理器必须实现的接口。
    """
    
    @abc.abstractmethod
    def load_document(self):
        """
        加载文档并将其转换为文本列表。
        
        Returns:
            list: 包含文档中每个段落的文本列表。
            如果加载失败，返回 None。
        """
        pass
    
    @abc.abstractmethod
    def process(self):
        """
        处理文档，包括加载和组合段落。
        
        Returns:
            list: 处理后的段落列表。
            如果处理失败，返回 None。
        """
        pass
    
    def combine_paragraphs(self, paragraphs):
        """
        根据配置的最大字符数合并段落。
        
        逻辑：逐个累加段落，直到累计字符数达到MAX_CHARS_PER_PARAGRAPH（默认3000），
        将累计到此的段落合并为一个新段落，添加到结果列表中，然后重新开始累计。
        
        Args:
            paragraphs (list): 原始段落列表。
            
        Returns:
            list: 合并后的段落列表。
        """
        if not paragraphs:
            return []
        
        combined_paragraphs = []
        current_combined = ""
        current_char_count = 0
        
        for paragraph in paragraphs:
            # 跳过空段落
            if not paragraph.strip():
                continue
                
            # 计算当前段落的字符数
            paragraph_chars = len(paragraph)
            
            # 判断是否需要开始新的组合段落
            if current_char_count + paragraph_chars > MAX_CHARS_PER_PARAGRAPH and current_combined:
                # 当前已累计的内容已经足够，添加到结果中
                combined_paragraphs.append(current_combined)
                # 重置累计变量
                current_combined = paragraph
                current_char_count = paragraph_chars
            else:
                # 继续累计
                if current_combined:
                    current_combined += " " + paragraph
                else:
                    current_combined = paragraph
                current_char_count += paragraph_chars
        
        # 添加最后一个累计的段落（如果有）
        if current_combined:
            combined_paragraphs.append(current_combined)
        
        return combined_paragraphs
    
    def preprocess_text(self, paragraphs):
        """
        对文本段落进行预处理，删除无意义的内容，使返回结果更整洁。
        
        Args:
            paragraphs (list): 需要处理的段落列表
            
        Returns:
            list: 处理后的段落列表
        """
        if not paragraphs:
            return []
            
        cleaned_paragraphs = []
        
        for paragraph in paragraphs:
            # 1. 去除空白字符
            paragraph = paragraph.strip()
            
            # 2. 跳过空段落
            if not paragraph:
                continue
                
            # 3. 跳过过短的段落（通常是页码、标点等）
            if len(paragraph) < 5:
                continue
                
            # 4. 替换多个空格为单个空格
            paragraph = re.sub(r'\s+', ' ', paragraph)
            
            # 5. 去除特殊字符和控制字符
            paragraph = re.sub(r'[\x00-\x1F\x7F-\x9F]', '', paragraph)
            
            # 6. 去除重复标点符号
            paragraph = re.sub(r'([。，！？；：,.!?;:])\1+', r'\1', paragraph)
            
            # 7. 去除纯数字行
            if re.match(r'^\d+(\s*\d+)*$', paragraph):
                continue
                
            # 8. 去除特殊格式行（如页眉页脚）
            if re.match(r'^[-=_*]{5,}$', paragraph):
                continue
                
            # 9. 去除URL和文件路径
            if re.match(r'^(https?://|www\.|file://|[a-zA-Z]:\\).+$', paragraph):
                continue
                
            # 10. 确保段落以句号或问号等结束
            if paragraph[-1] not in '.。!！?？;；':
                paragraph += '。'
                
            # 添加到清洁段落列表
            cleaned_paragraphs.append(paragraph)
            
        # 过滤掉重复段落
        unique_paragraphs = []
        for paragraph in cleaned_paragraphs:
            is_duplicate = False
            for existing in unique_paragraphs:
                # 检查是否是子字符串或相似度高
                if paragraph in existing or existing in paragraph:
                    is_duplicate = True
                    break
            if not is_duplicate:
                unique_paragraphs.append(paragraph)
                
        return unique_paragraphs

class PDFDocumentProcessor(DocumentProcessor):
    """
    处理 PDF 文档的类，将其转换为文本并根据规则组合段落。
    
    Attributes:
        input_path (str): 输入的 PDF 文档路径。
    """
    
    def __init__(self, input_path):
        """
        初始化 PDFDocumentProcessor 类的实例。
        
        Args:
            input_path (str): 输入的 PDF 文档路径。
        """
        if not isinstance(input_path, str):
            raise TypeError("input_path must be a string")
        self.input_path = os.path.abspath(input_path)  # 确保使用绝对路径
        
    def load_document(self):
        """
        加载 PDF 文档并将其转换为文本列表。
        
        Returns:
            list: 包含文档中每个段落的文本列表。
            如果加载失败，返回 None。
        """
        try:
            paragraphs = []
            with open(self.input_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    text = page.extract_text()
                    
                    # 按行分割文本
                    lines = text.split('\n')
                    
                    # 处理每一行
                    for line in lines:
                        clean_line = line.strip()
                        if clean_line:  # 只添加非空行
                            paragraphs.append(clean_line)
            
            return paragraphs if paragraphs else None
        except Exception as e:
            print(f"加载PDF文档失败: {e}")
            return None
    
    def combine_paragraphs(self, paragraphs):
        """
        根据配置的最大字符数合并段落。
        
        Args:
            paragraphs (list): 原始段落列表。
            
        Returns:
            list: 合并后的段落列表。
        """
        # 使用基类的实现
        return super().combine_paragraphs(paragraphs)
    
    def process(self):
        """
        处理 PDF 文档，包括加载和组合段落。
        
        Returns:
            list: 处理后的段落列表。
            如果处理失败，返回 None。
        """
        paragraphs = self.load_document()
        if paragraphs is None:
            return None
        
        combined_paragraphs = self.combine_paragraphs(paragraphs)
        # 应用文本预处理
        return self.preprocess_text(combined_paragraphs)

class TXTDocumentProcessor(DocumentProcessor):
    """
    处理纯文本文档的类，将其转换为文本并根据规则组合段落。
    
    Attributes:
        input_path (str): 输入的文本文档路径。
    """
    
    def __init__(self, input_path):
        """
        初始化 TXTDocumentProcessor 类的实例。
        
        Args:
            input_path (str): 输入的文本文档路径。
        """
        if not isinstance(input_path, str):
            raise TypeError("input_path must be a string")
        self.input_path = os.path.abspath(input_path)  # 确保使用绝对路径
    
    def _detect_encoding(self, file_path):
        """
        检测文件编码。
        
        Args:
            file_path (str): 文件路径。
            
        Returns:
            str: 文件编码。
        """
        with open(file_path, 'rb') as f:
            result = chardet.detect(f.read())
        return result['encoding'] or 'utf-8'  # 默认使用 utf-8
    
    def load_document(self):
        """
        加载文本文档并将其转换为文本列表。
        
        Returns:
            list: 包含文档中每个段落的文本列表。
            如果加载失败，返回 None。
        """
        try:
            encoding = self._detect_encoding(self.input_path)
            with open(self.input_path, 'r', encoding=encoding, errors='replace') as file:
                content = file.read()
            
            # 按空行分割段落
            paragraphs = []
            for para in re.split(r'\n\s*\n', content):
                clean_para = para.strip()
                if clean_para:  # 只添加非空段落
                    # 替换段落内的换行符为空格
                    clean_para = re.sub(r'\n', ' ', clean_para)
                    paragraphs.append(clean_para)
            
            return paragraphs if paragraphs else None
        except Exception as e:
            print(f"加载文本文档失败: {e}")
            return None
    
    def combine_paragraphs(self, paragraphs):
        """
        根据配置的最大字符数合并段落。
        
        Args:
            paragraphs (list): 原始段落列表。
            
        Returns:
            list: 合并后的段落列表。
        """
        # 使用基类的实现
        return super().combine_paragraphs(paragraphs)
    
    def process(self):
        """
        处理文本文档，包括加载和组合段落。
        
        Returns:
            list: 处理后的段落列表。
            如果处理失败，返回 None。
        """
        paragraphs = self.load_document()
        if paragraphs is None:
            return None
        
        combined_paragraphs = self.combine_paragraphs(paragraphs)
        # 应用文本预处理
        return self.preprocess_text(combined_paragraphs)

class CSVDocumentProcessor(DocumentProcessor):
    """
    处理 CSV 文档的类，将其转换为文本并根据规则组合段落。
    
    Attributes:
        input_path (str): 输入的 CSV 文档路径。
    """
    
    def __init__(self, input_path):
        """
        初始化 CSVDocumentProcessor 类的实例。
        
        Args:
            input_path (str): 输入的 CSV 文档路径。
        """
        if not isinstance(input_path, str):
            raise TypeError("input_path must be a string")
        self.input_path = os.path.abspath(input_path)  # 确保使用绝对路径
    
    def _detect_encoding(self, file_path):
        """
        检测文件编码。
        
        Args:
            file_path (str): 文件路径。
            
        Returns:
            str: 文件编码。
        """
        with open(file_path, 'rb') as f:
            result = chardet.detect(f.read())
        return result['encoding'] or 'utf-8'  # 默认使用 utf-8
    
    def load_document(self):
        """
        加载 CSV 文档并将其转换为文本列表。
        
        Returns:
            list: 包含文档中每个行的文本列表。
            如果加载失败，返回 None。
        """
        try:
            encoding = self._detect_encoding(self.input_path)
            paragraphs = []
            
            with open(self.input_path, 'r', encoding=encoding, errors='replace') as file:
                csv_reader = csv.reader(file)
                headers = next(csv_reader, None)  # 获取表头
                
                if headers:
                    paragraphs.append(', '.join(headers))  # 添加表头作为第一段
                
                for row in csv_reader:
                    if row:  # 只添加非空行
                        paragraphs.append(', '.join(row))
            
            return paragraphs if paragraphs else None
        except Exception as e:
            print(f"加载CSV文档失败: {e}")
            return None
    
    def combine_paragraphs(self, paragraphs):
        """
        根据配置的最大字符数合并段落。
        
        Args:
            paragraphs (list): 原始段落列表。
            
        Returns:
            list: 合并后的段落列表。
        """
        # 使用基类的实现
        return super().combine_paragraphs(paragraphs)
    
    def process(self):
        """
        处理 CSV 文档，包括加载和组合段落。
        
        Returns:
            list: 处理后的段落列表。
            如果处理失败，返回 None。
        """
        paragraphs = self.load_document()
        if paragraphs is None:
            return None
        
        combined_paragraphs = self.combine_paragraphs(paragraphs)
        # 应用文本预处理
        return self.preprocess_text(combined_paragraphs)

class ExcelDocumentProcessor(DocumentProcessor):
    """
    处理 Excel 文档的类，将其转换为文本并根据规则组合段落。
    
    Attributes:
        input_path (str): 输入的 Excel 文档路径。
    """
    
    def __init__(self, input_path):
        """
        初始化 ExcelDocumentProcessor 类的实例。
        
        Args:
            input_path (str): 输入的 Excel 文档路径。
        """
        if not isinstance(input_path, str):
            raise TypeError("input_path must be a string")
        self.input_path = os.path.abspath(input_path)  # 确保使用绝对路径
    
    def load_document(self):
        """
        加载 Excel 文档并将其转换为文本列表。
        
        Returns:
            list: 包含文档中每个工作表和单元格的文本列表。
            如果加载失败，返回 None。
        """
        try:
            workbook = openpyxl.load_workbook(self.input_path, data_only=True)
            paragraphs = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                
                # 添加工作表名称作为段落
                paragraphs.append(f"工作表: {sheet_name}")
                
                # 获取所有行
                rows = list(sheet.rows)
                if not rows:
                    continue
                
                # 获取表头
                headers = [cell.value for cell in rows[0]]
                
                # 处理每一行数据
                for row in rows[1:]:
                    row_data = []
                    for i, cell in enumerate(row):
                        if cell.value is not None:
                            if i < len(headers) and headers[i] is not None:
                                row_data.append(f"{headers[i]}: {cell.value}")
                            else:
                                row_data.append(str(cell.value))
                    
                    if row_data:  # 只添加非空行
                        paragraphs.append(', '.join(row_data))
            
            return paragraphs if paragraphs else None
        except Exception as e:
            print(f"加载Excel文档失败: {e}")
            return None
    
    def combine_paragraphs(self, paragraphs):
        """
        根据配置的最大字符数合并段落。
        
        Args:
            paragraphs (list): 原始段落列表。
            
        Returns:
            list: 合并后的段落列表。
        """
        # 使用基类的实现
        return super().combine_paragraphs(paragraphs)
    
    def process(self):
        """
        处理 Excel 文档，包括加载和组合段落。
        
        Returns:
            list: 处理后的段落列表。
            如果处理失败，返回 None。
        """
        paragraphs = self.load_document()
        if paragraphs is None:
            return None
        
        combined_paragraphs = self.combine_paragraphs(paragraphs)
        # 应用文本预处理
        return self.preprocess_text(combined_paragraphs)

class MarkdownDocumentProcessor(DocumentProcessor):
    """
    处理 Markdown 文档的类，将其转换为文本并根据规则组合段落。
    
    Attributes:
        input_path (str): 输入的 Markdown 文档路径。
    """
    
    def __init__(self, input_path):
        """
        初始化 MarkdownDocumentProcessor 类的实例。
        
        Args:
            input_path (str): 输入的 Markdown 文档路径。
        """
        if not isinstance(input_path, str):
            raise TypeError("input_path must be a string")
        self.input_path = os.path.abspath(input_path)  # 确保使用绝对路径
    
    def _detect_encoding(self, file_path):
        """
        检测文件编码。
        
        Args:
            file_path (str): 文件路径。
            
        Returns:
            str: 文件编码。
        """
        with open(file_path, 'rb') as f:
            result = chardet.detect(f.read())
        return result['encoding'] or 'utf-8'  # 默认使用 utf-8
    
    def load_document(self):
        """
        加载 Markdown 文档并将其转换为文本列表。
        
        Returns:
            list: 包含文档中每个段落的文本列表。
            如果加载失败，返回 None。
        """
        try:
            encoding = self._detect_encoding(self.input_path)
            with open(self.input_path, 'r', encoding=encoding, errors='replace') as file:
                content = file.read()
            
            # 将 Markdown 转换为 HTML
            html = markdown.markdown(content)
            
            # 将 HTML 转换为纯文本
            h = html2text.HTML2Text()
            h.ignore_links = True
            h.ignore_images = True
            h.ignore_emphasis = True
            text = h.handle(html)
            
            # 按空行分割段落
            paragraphs = []
            for para in re.split(r'\n\s*\n', text):
                clean_para = para.strip()
                if clean_para:  # 只添加非空段落
                    # 替换段落内的换行符为空格
                    clean_para = re.sub(r'\n', ' ', clean_para)
                    paragraphs.append(clean_para)
            
            return paragraphs if paragraphs else None
        except Exception as e:
            print(f"加载Markdown文档失败: {e}")
            return None
    
    def combine_paragraphs(self, paragraphs):
        """
        根据配置的最大字符数合并段落。
        
        Args:
            paragraphs (list): 原始段落列表。
            
        Returns:
            list: 合并后的段落列表。
        """
        # 使用基类的实现
        return super().combine_paragraphs(paragraphs)
    
    def process(self):
        """
        处理 Markdown 文档，包括加载和组合段落。
        
        Returns:
            list: 处理后的段落列表。
            如果处理失败，返回 None。
        """
        paragraphs = self.load_document()
        if paragraphs is None:
            return None
        
        combined_paragraphs = self.combine_paragraphs(paragraphs)
        # 应用文本预处理
        return self.preprocess_text(combined_paragraphs)

class PowerPointDocumentProcessor(DocumentProcessor):
    """
    处理 PowerPoint 文档的类，将其转换为文本并根据规则组合段落。
    
    Attributes:
        input_path (str): 输入的 PowerPoint 文档路径。
    """
    
    def __init__(self, input_path):
        """
        初始化 PowerPointDocumentProcessor 类的实例。
        
        Args:
            input_path (str): 输入的 PowerPoint 文档路径。
        """
        if not isinstance(input_path, str):
            raise TypeError("input_path must be a string")
        self.input_path = os.path.abspath(input_path)  # 确保使用绝对路径
    
    def load_document(self):
        """
        加载 PowerPoint 文档并将其转换为文本列表。
        
        Returns:
            list: 包含文档中每个幻灯片的文本列表。
            如果加载失败，返回 None。
        """
        try:
            presentation = pptx.Presentation(self.input_path)
            paragraphs = []
            
            for i, slide in enumerate(presentation.slides):
                slide_text = []
                slide_text.append(f"幻灯片 {i+1}")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text.strip())
                
                if slide_text:  # 只添加非空幻灯片
                    paragraphs.append("\n".join(slide_text))
            
            return paragraphs if paragraphs else None
        except Exception as e:
            print(f"加载PowerPoint文档失败: {e}")
            return None
    
    def combine_paragraphs(self, paragraphs):
        """
        根据配置的最大字符数合并段落。
        
        Args:
            paragraphs (list): 原始段落列表。
            
        Returns:
            list: 合并后的段落列表。
        """
        # 使用基类的实现
        return super().combine_paragraphs(paragraphs)
    
    def process(self):
        """
        处理 PowerPoint 文档，包括加载和组合段落。
        
        Returns:
            list: 处理后的段落列表。
            如果处理失败，返回 None。
        """
        paragraphs = self.load_document()
        if paragraphs is None:
            return None
        
        combined_paragraphs = self.combine_paragraphs(paragraphs)
        # 应用文本预处理
        return self.preprocess_text(combined_paragraphs)

class WordDocumentProcessor(DocumentProcessor):
    """
    处理 Word 文档的类，将其转换为文本并根据规则组合段落。
    支持.doc和.docx格式。跨平台实现，兼容Windows和Linux。

    Attributes:
        input_docx_path (str): 输入的 Word 文档路径。
    """

    def __init__(self, input_docx_path):
        """
        初始化 WordDocumentProcessor 类的实例。

        Args:
            input_docx_path (str): 输入的 Word 文档路径。
        """
        if not isinstance(input_docx_path, str):
            raise TypeError("input_docx_path must be a string")
        self.input_docx_path = os.path.abspath(input_docx_path)  # 确保使用绝对路径
        self.is_doc_format = self._is_doc_format()
        self.temp_files = []  # 跟踪所有创建的临时文件

    def __del__(self):
        """
        析构函数，确保所有临时文件被清理。
        """
        self.cleanup_temp_files()

    def cleanup_temp_files(self):
        """
        清理所有临时文件。
        """
        for file_path in self.temp_files:
            try:
                if os.path.exists(file_path):
                    os.remove(file_path)
                    print(f"已删除临时文件: {file_path}")
            except Exception as e:
                print(f"删除临时文件失败 {file_path}: {e}")
        self.temp_files = []

    def _is_doc_format(self):
        """
        检查文件是否为旧版.doc格式。

        Returns:
            bool: 如果是.doc格式返回True，否则返回False。
        """
        _, ext = os.path.splitext(self.input_docx_path)
        return ext.lower() == '.doc'

    def _convert_doc_to_docx_windows(self):
        """
        Windows平台使用COM组件将.doc格式转换为.docx格式。

        Returns:
            str: 转换后的.docx文件路径，如果转换失败则返回None。
        """
        if not IS_WINDOWS:
            return None
            
        try:
            # 初始化COM组件
            pythoncom.CoInitialize()
            
            # 创建临时文件路径 - 使用更简单的文件名和路径
            temp_dir = os.path.join(tempfile.gettempdir(), "word_processor")
            os.makedirs(temp_dir, exist_ok=True)
            
            # 复制原文件到临时目录，使用简单文件名
            temp_filename = f"temp_{int(time.time())}.doc"
            temp_doc_path = os.path.join(temp_dir, temp_filename)
            shutil.copy2(self.input_docx_path, temp_doc_path)
            self.temp_files.append(temp_doc_path)  # 添加到临时文件列表
            
            # 设置输出文件路径
            docx_path = os.path.join(temp_dir, f"{os.path.splitext(temp_filename)[0]}.docx")
            self.temp_files.append(docx_path)  # 添加到临时文件列表
            
            print(f"原始文件路径: {self.input_docx_path}")
            print(f"临时文件路径: {temp_doc_path}")
            print(f"目标文件路径: {docx_path}")
            
            # 使用Word COM对象转换文件
            word = wc.Dispatch('Word.Application')
            word.Visible = False
            
            # 尝试打开文档
            try:
                doc = word.Documents.Open(temp_doc_path)
                # 保存为docx格式
                doc.SaveAs2(docx_path, FileFormat=16)  # 16表示docx格式
                doc.Close()
            except Exception as e:
                print(f"打开或保存文档时出错: {e}")
                # 尝试使用另一种方式
                if os.path.exists(temp_doc_path):
                    try:
                        # 使用shell命令转换
                        print("尝试使用替代方法转换文档...")
                        word.Quit()
                        time.sleep(1)  # 确保Word完全关闭
                        
                        # 使用python-docx直接尝试打开.doc文件
                        # 有些.doc文件实际上是可以被python-docx打开的
                        try:
                            doc = Document(temp_doc_path)
                            doc.save(docx_path)
                            print("使用python-docx成功转换文档")
                            return docx_path
                        except Exception as e2:
                            print(f"python-docx无法打开文档: {e2}")
                            return None
                    except Exception as e3:
                        print(f"替代转换方法失败: {e3}")
                        return None
            finally:
                try:
                    word.Quit()
                except:
                    pass
            
            # 检查文件是否成功创建
            if os.path.exists(docx_path):
                return docx_path
            return None
        except Exception as e:
            print(f"Windows平台转换.doc到.docx失败: {e}")
            return None
        finally:
            # 释放COM组件
            pythoncom.CoUninitialize()

    def _convert_doc_to_docx_linux(self):
        """
        Linux平台使用LibreOffice将.doc格式转换为.docx格式。

        Returns:
            str: 转换后的.docx文件路径，如果转换失败则返回None。
        """
        try:
            # 创建临时文件路径
            temp_dir = os.path.join(tempfile.gettempdir(), "word_processor")
            os.makedirs(temp_dir, exist_ok=True)
            
            # 复制原文件到临时目录，使用简单文件名
            temp_filename = f"temp_{int(time.time())}.doc"
            temp_doc_path = os.path.join(temp_dir, temp_filename)
            shutil.copy2(self.input_docx_path, temp_doc_path)
            self.temp_files.append(temp_doc_path)  # 添加到临时文件列表
            
            # 设置输出文件路径
            docx_path = os.path.join(temp_dir, f"{os.path.splitext(temp_filename)[0]}.docx")
            self.temp_files.append(docx_path)  # 添加到临时文件列表
            
            print(f"原始文件路径: {self.input_docx_path}")
            print(f"临时文件路径: {temp_doc_path}")
            print(f"目标文件路径: {docx_path}")
            
            # 检查是否安装了LibreOffice或OpenOffice
            libreoffice_cmd = shutil.which('libreoffice')
            soffice_cmd = shutil.which('soffice')
            
            if libreoffice_cmd:
                cmd = [libreoffice_cmd, '--headless', '--convert-to', 'docx', '--outdir', temp_dir, temp_doc_path]
                print(f"使用LibreOffice转换文档: {' '.join(cmd)}")
                result = subprocess.run(cmd, capture_output=True, text=True)
                if result.returncode != 0:
                    print(f"LibreOffice转换失败: {result.stderr}")
                    return None
            elif soffice_cmd:
                cmd = [soffice_cmd, '--headless', '--convert-to', 'docx', '--outdir', temp_dir, temp_doc_path]
                print(f"使用OpenOffice转换文档: {' '.join(cmd)}")
                result = subprocess.run(cmd, capture_output=True, text=True)
                if result.returncode != 0:
                    print(f"OpenOffice转换失败: {result.stderr}")
                    return None
            else:
                # 如果没有安装LibreOffice或OpenOffice，尝试使用其他工具
                print("未找到LibreOffice或OpenOffice，尝试使用其他方法...")
                
                # 尝试使用unoconv (如果已安装)
                if shutil.which('unoconv'):
                    cmd = ['unoconv', '-f', 'docx', '-o', docx_path, temp_doc_path]
                    print(f"使用unoconv转换文档: {' '.join(cmd)}")
                    result = subprocess.run(cmd, capture_output=True, text=True)
                    if result.returncode != 0:
                        print(f"unoconv转换失败: {result.stderr}")
                        return None
                else:
                    print("未找到合适的文档转换工具")
                    return None
            
            # 尝试找到生成的文件 - LibreOffice可能会使用不同的文件名
            if not os.path.exists(docx_path):
                # 查找可能的输出文件
                potential_docx_files = [f for f in os.listdir(temp_dir) if f.endswith('.docx') and f.startswith(os.path.splitext(temp_filename)[0])]
                if potential_docx_files:
                    # 使用最新创建的文件
                    newest_file = max(potential_docx_files, key=lambda f: os.path.getmtime(os.path.join(temp_dir, f)))
                    docx_path = os.path.join(temp_dir, newest_file)
                    self.temp_files.append(docx_path)
                    print(f"找到转换后的文件: {docx_path}")
            
            # 检查文件是否成功创建
            if os.path.exists(docx_path):
                return docx_path
                
            return None
        except Exception as e:
            print(f"Linux平台转换.doc到.docx失败: {e}")
            return None

    def _convert_doc_to_docx(self):
        """
        跨平台实现: 将.doc格式转换为.docx格式。
        根据系统平台选择适当的转换方法。

        Returns:
            str: 转换后的.docx文件路径，如果转换失败则返回None。
        """
        # 根据平台选择不同的实现
        if IS_WINDOWS:
            return self._convert_doc_to_docx_windows()
        else:
            return self._convert_doc_to_docx_linux()

    def _extract_text_directly(self):
        """
        直接从文件中提取文本，不依赖外部库。
        
        Returns:
            list: 包含文档中每个段落的文本列表。
            如果提取失败，返回None。
        """
        try:
            print("尝试直接提取文本...")
            with open(self.input_docx_path, 'rb') as f:
                content = f.read()
                
            # 方法1: 尝试提取ASCII文本
            text = ""
            for i in range(0, len(content)):
                if 32 <= content[i] <= 126 or content[i] in (10, 13):  # ASCII可打印字符或换行符
                    text += chr(content[i])
            
            # 清理文本
            text = re.sub(r'[^\w\s\.,;:!?\'"-]', ' ', text)  # 保留基本标点和字母数字
            text = re.sub(r'\s+', ' ', text)  # 合并多个空格
            
            # 分割成段落 - 使用多种可能的段落分隔符
            paragraphs = []
            for para in re.split(r'\n\s*\n|\r\n\s*\r\n', text):
                clean_para = para.strip()
                if clean_para and len(clean_para) > 10:  # 只保留有意义的段落
                    paragraphs.append(clean_para)
            
            if paragraphs:
                print(f"直接提取成功，找到{len(paragraphs)}个段落")
                return paragraphs
                
            # 方法2: 尝试提取Unicode文本
            text = ""
            for i in range(0, len(content), 2):
                if i+1 < len(content):
                    char_code = content[i] + (content[i+1] << 8)
                    if 32 <= char_code <= 126 or char_code in (10, 13):
                        text += chr(char_code)
            
            # 清理和分割文本
            paragraphs = []
            for para in re.split(r'\n\s*\n|\r\n\s*\r\n', text):
                clean_para = para.strip()
                if clean_para and len(clean_para) > 10:
                    paragraphs.append(clean_para)
            
            if paragraphs:
                print(f"Unicode提取成功，找到{len(paragraphs)}个段落")
                return paragraphs
            
            return None
        except Exception as e:
            print(f"直接提取文本失败: {e}")
            return None

    def load_document(self):
        """
        加载 Word 文档并将其转换为文本列表。
        支持.doc和.docx格式。

        Returns:
            list: 包含文档中每个段落的文本列表。
            如果加载失败，返回 None。
        """
        try:
            # 如果是.docx格式，直接使用python-docx打开
            if not self.is_doc_format:
                try:
                    doc = Document(self.input_docx_path)
                    full_text = []
                    for paragraph in doc.paragraphs:
                        clean_text = paragraph.text.strip()
                        if clean_text:
                            full_text.append(clean_text)
                    return full_text
                except Exception as e:
                    print(f"打开.docx文件失败: {e}")
                    # 尝试使用备用方法
                    return self._fallback_load_document()
            
            # 如果是.doc格式，先转换为.docx
            print(f"检测到.doc格式文件，尝试转换为.docx...")
            temp_file = self._convert_doc_to_docx()
            if temp_file:
                try:
                    # 使用python-docx加载文档
                    doc = Document(temp_file)
                    full_text = []
                    for paragraph in doc.paragraphs:
                        # 去除每个段落开头和结尾的空白字符
                        clean_text = paragraph.text.strip()
                        if clean_text:  # 仅在段落不为空时添加
                            full_text.append(clean_text)
                    
                    return full_text
                except Exception as e:
                    print(f"打开转换后的.docx文件失败: {e}")
            
            # 如果转换失败，尝试直接提取文本
            text_from_direct = self._extract_text_directly()
            if text_from_direct:
                return text_from_direct
                
            # 如果直接提取也失败，尝试使用备用方法
            return self._fallback_load_document()
        except Exception as e:
            print(f"Error loading document: {e}")
            # 尝试使用备用方法读取文件内容
            return self._fallback_load_document()

    def _fallback_load_document(self):
        """
        当主要加载方法失败时的备用方法。
        尝试使用各种方式读取文档内容，支持多平台。

        Returns:
            list: 包含文档中每个段落的文本列表。
            如果加载失败，返回 None。
        """
        try:
            print("尝试使用备用方法读取文档...")
            
            # 1. 尝试使用简单的文本提取方法 (跨平台)
            try:
                with open(self.input_docx_path, 'rb') as f:
                    content = f.read()
                    # 尝试提取可读文本
                    text = ""
                    for i in range(0, len(content), 2):
                        if i+1 < len(content):
                            char = content[i:i+2]
                            if char[0] >= 32 and char[0] <= 126 and char[1] == 0:
                                text += chr(char[0])
                    
                    # 分割成段落
                    paragraphs = [p.strip() for p in text.split('\n') if p.strip()]
                    if paragraphs and len(paragraphs) > 5:  # 确保提取到足够的段落
                        print("使用简单文本提取成功")
                        return paragraphs
            except Exception as e:
                print(f"简单文本提取失败: {e}")
            
            # 2. 尝试使用跨平台工具 (如果安装了)
            
            # 2.1 antiword (适用于Linux)
            try:
                if shutil.which('antiword'):
                    print("尝试使用antiword提取文本...")
                    result = subprocess.run(['antiword', self.input_docx_path], 
                                           capture_output=True, text=True, check=True)
                    if result.stdout:
                        paragraphs = [p.strip() for p in result.stdout.split('\n') if p.strip()]
                        if paragraphs and len(paragraphs) > 5:
                            print("使用antiword提取成功")
                            return paragraphs
            except Exception as e:
                print(f"antiword提取失败: {e}")
            
            # 2.2 catdoc (适用于Linux)
            try:
                if shutil.which('catdoc'):
                    print("尝试使用catdoc提取文本...")
                    result = subprocess.run(['catdoc', self.input_docx_path], 
                                           capture_output=True, text=True, check=True)
                    if result.stdout:
                        paragraphs = [p.strip() for p in result.stdout.split('\n') if p.strip()]
                        if paragraphs and len(paragraphs) > 5:
                            print("使用catdoc提取成功")
                            return paragraphs
            except Exception as e:
                print(f"catdoc提取失败: {e}")
            
            # 3. 平台特定的方法
            
            # 3.1 Windows COM方法 (仅适用于Windows)
            if IS_WINDOWS:
                try:
                    # 初始化COM组件
                    pythoncom.CoInitialize()
                    
                    # 复制文件到临时位置
                    temp_dir = os.path.join(tempfile.gettempdir(), "word_processor")
                    os.makedirs(temp_dir, exist_ok=True)
                    temp_filename = f"fallback_{int(time.time())}"
                    temp_doc_path = os.path.join(temp_dir, f"{temp_filename}.doc")
                    shutil.copy2(self.input_docx_path, temp_doc_path)
                    self.temp_files.append(temp_doc_path)  # 添加到临时文件列表
                    
                    # 使用Word COM对象直接读取文本
                    word = wc.Dispatch('Word.Application')
                    word.Visible = False
                    
                    try:
                        doc = word.Documents.Open(temp_doc_path)
                        
                        full_text = []
                        for i in range(1, doc.Paragraphs.Count + 1):
                            text = doc.Paragraphs(i).Range.Text.strip()
                            if text:
                                full_text.append(text)
                        
                        doc.Close()
                        
                        # 释放COM组件
                        pythoncom.CoUninitialize()
                        
                        if full_text and len(full_text) > 5:
                            print("使用Word COM直接读取成功")
                            return full_text
                    except Exception as e:
                        print(f"备用方法使用Word COM失败: {e}")
                        
                        # 最后尝试：使用Word的SaveAs方法保存为纯文本
                        try:
                            print("尝试将文档保存为纯文本...")
                            doc = word.Documents.Open(temp_doc_path)
                            txt_path = os.path.join(temp_dir, f"{os.path.splitext(temp_filename)[0]}.txt")
                            self.temp_files.append(txt_path)  # 添加到临时文件列表
                            doc.SaveAs(txt_path, FileFormat=2)  # 2表示纯文本格式
                            doc.Close()
                            
                            # 读取保存的文本文件
                            with open(txt_path, 'r', encoding='utf-8', errors='ignore') as f:
                                lines = f.readlines()
                            
                            # 处理文本行
                            paragraphs = []
                            current_para = ""
                            for line in lines:
                                line = line.strip()
                                if not line:  # 空行表示段落分隔
                                    if current_para:
                                        paragraphs.append(current_para)
                                        current_para = ""
                                else:
                                    if current_para:
                                        current_para += " " + line
                                    else:
                                        current_para = line
                            
                            # 添加最后一个段落
                            if current_para:
                                paragraphs.append(current_para)
                            
                            if paragraphs and len(paragraphs) > 5:
                                print("通过保存为文本文件成功提取内容")
                                return paragraphs
                        except Exception as e:
                            print(f"保存为文本文件失败: {e}")
                    finally:
                        try:
                            word.Quit()
                        except:
                            pass
                        
                        # 释放COM组件 (如果之前未释放)
                        try:
                            pythoncom.CoUninitialize()
                        except:
                            pass
                except Exception as e:
                    print(f"Windows COM方法失败: {e}")
            
            # 3.2 Linux特定方法
            else:
                # 尝试使用LibreOffice/OpenOffice转换为文本
                try:
                    temp_dir = os.path.join(tempfile.gettempdir(), "word_processor")
                    os.makedirs(temp_dir, exist_ok=True)
                    temp_filename = f"fallback_{int(time.time())}"
                    temp_doc_path = os.path.join(temp_dir, f"{temp_filename}.doc")
                    shutil.copy2(self.input_docx_path, temp_doc_path)
                    self.temp_files.append(temp_doc_path)
                    
                    txt_path = os.path.join(temp_dir, f"{temp_filename}.txt")
                    self.temp_files.append(txt_path)
                    
                    libreoffice_cmd = shutil.which('libreoffice')
                    soffice_cmd = shutil.which('soffice')
                    
                    print("尝试使用LibreOffice/OpenOffice转换为文本...")
                    
                    if libreoffice_cmd:
                        cmd = [libreoffice_cmd, '--headless', '--convert-to', 'txt', '--outdir', temp_dir, temp_doc_path]
                        subprocess.run(cmd, capture_output=True, text=True, check=True)
                    elif soffice_cmd:
                        cmd = [soffice_cmd, '--headless', '--convert-to', 'txt', '--outdir', temp_dir, temp_doc_path]
                        subprocess.run(cmd, capture_output=True, text=True, check=True)
                    
                    # 查找转换后的文本文件
                    for file in os.listdir(temp_dir):
                        if file.startswith(temp_filename) and file.endswith('.txt'):
                            txt_path = os.path.join(temp_dir, file)
                            break
                    
                    if os.path.exists(txt_path):
                        with open(txt_path, 'r', encoding='utf-8', errors='ignore') as f:
                            lines = f.readlines()
                        
                        # 处理文本行
                        paragraphs = []
                        current_para = ""
                        for line in lines:
                            line = line.strip()
                            if not line:  # 空行表示段落分隔
                                if current_para:
                                    paragraphs.append(current_para)
                                    current_para = ""
                            else:
                                if current_para:
                                    current_para += " " + line
                                else:
                                    current_para = line
                        
                        # 添加最后一个段落
                        if current_para:
                            paragraphs.append(current_para)
                        
                        if paragraphs and len(paragraphs) > 5:
                            print("通过LibreOffice/OpenOffice转换为文本成功")
                            return paragraphs
                except Exception as e:
                    print(f"LibreOffice/OpenOffice转换失败: {e}")
            
            # 4. 最后尝试: 使用textract库 (如果已安装)
            try:
                import importlib
                if importlib.util.find_spec("textract"):
                    import textract
                    print("尝试使用textract库提取文本...")
                    text = textract.process(self.input_docx_path).decode('utf-8', errors='ignore')
                    if text:
                        # 分割成段落
                        paragraphs = [p.strip() for p in text.split('\n') if p.strip()]
                        if paragraphs and len(paragraphs) > 5:
                            print("使用textract提取成功")
                            return paragraphs
            except Exception as e:
                print(f"textract提取失败: {e}")
                
            print("所有备用方法均失败")
            return None
        except Exception as e:
            print(f"备用方法也失败了: {e}")
            return None

    def combine_paragraphs(self, paragraphs):
        """
        根据配置的最大字符数合并段落。

        Args:
            paragraphs (list): 原始段落列表。

        Returns:
            list: 合并后的段落列表。
        """
        # 使用基类的实现
        return super().combine_paragraphs(paragraphs)

    def process(self):
        """
        处理 Word 文档，包括加载和组合段落。

        Returns:
            list: 处理后的段落列表。
            如果处理失败，返回 None。
        """
        paragraphs = self.load_document()
        if paragraphs is None:
            return None
            
        combined_paragraphs = self.combine_paragraphs(paragraphs)
        # 应用文本预处理
        return self.preprocess_text(combined_paragraphs)

if __name__ == "__main__":
    input_docx_path = "test/input.docx"
    processor = WordDocumentProcessor(input_docx_path)
    combined_paragraphs = processor.process()
    if combined_paragraphs:
        for i, paragraph in enumerate(combined_paragraphs):
            print(f"段落 {i+1}: {paragraph[:50]}...")  # 打印每个大段的前50个字符

class DocumentProcessorFactory:
    """
    文档处理器工厂类，用于根据文件扩展名创建相应的文档处理器。
    """
    
    @staticmethod
    def create_processor(file_path):
        """
        根据文件扩展名创建相应的文档处理器。
        
        Args:
            file_path (str): 文件路径。
            
        Returns:
            DocumentProcessor: 相应的文档处理器实例。
            如果不支持该文件类型，则返回 None。
        """
        if not isinstance(file_path, str):
            raise TypeError("file_path must be a string")
        
        _, ext = os.path.splitext(file_path)
        ext = ext.lower()
        
        # 根据文件扩展名创建相应的处理器
        if ext in ['.doc', '.docx', '.docm', '.dotm']:
            return WordDocumentProcessor(file_path)
        elif ext == '.pdf':
            return PDFDocumentProcessor(file_path)
        elif ext in ['.txt', '.text']:
            return TXTDocumentProcessor(file_path)
        elif ext == '.csv':
            return CSVDocumentProcessor(file_path)
        elif ext in ['.xls', '.xlsx', '.xlsm']:
            return ExcelDocumentProcessor(file_path)
        elif ext in ['.md', '.markdown']:
            return MarkdownDocumentProcessor(file_path)
        elif ext in ['.ppt', '.pptx', '.pptm']:
            return PowerPointDocumentProcessor(file_path)
        else:
            print(f"不支持的文件类型: {ext}")
            return None

